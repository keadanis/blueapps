# -*- coding: utf-8 -*-
#############################################################################
#
#    Copyright (C) 2019-Antti Kärki.
#    Author: Antti Kärki.
#
#    You can modify it under the terms of the GNU AFFERO
#    GENERAL PUBLIC LICENSE (AGPL v3), Version 3.
#
#    This program is distributed in the hope that it will be useful,
#    but WITHOUT ANY WARRANTY; without even the implied warranty of
#    MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
#    GNU AFFERO GENERAL PUBLIC LICENSE (AGPL v3) for more details.
#
#    You should have received a copy of the GNU AFFERO GENERAL PUBLIC LICENSE
#    (AGPL v3) along with this program.
#    If not, see <http://www.gnu.org/licenses/>.
#
# pip3 install python-pptx
# pip3 install openpyxl
# pip3 install pandas
#
#############################################################################

from odoo import api, fields, models
from odoo import exceptions
from odoo import http
import logging
import tempfile
from datetime import date, datetime, timedelta
import base64
from . import rocker_connection
from . import rocker_about
import openpyxl as op
from openpyxl.styles import Font, Fill, Alignment
import pandas as pd
from openpyxl.worksheet.table import Table, TableStyleInfo
from openpyxl.utils import get_column_letter
import io



_logger = logging.getLogger(__name__)
# _logger.debug('A DEBUG message')
# _logger.info('An INFO message')
# _logger.warning('A WARNING message')
# _logger.error('An ERROR message')
# --log_level=:DEBUG

class rocker_report_collection(models.Model):
    _name = 'rocker.report.collection'
    _description = 'Rocker Collection Reports'
    _order = 'sequence'
    report_id = fields.Integer('Report_id')
    collection_id = fields.Integer('Collection_id')
    name = fields.Char('Name')
    sequence = fields.Integer('sequence', help="Sequence for the handle.",default=10)



class Report(models.Model):

    _name = 'rocker.report'
    _description = 'Rocker Reporting'
    name = fields.Char('Name', required=True)
    report_type = fields.Selection(
        [('single', 'Single'), ('collection', 'Collection')], 'Report type', default='single')
    description = fields.Text('Description')
    active = fields.Boolean('Active', default=True)
    schedule_onoff = fields.Boolean('Scheduling Active', default=False)
    store_history = fields.Boolean('Store to Archive', default=False)
    date_published = fields.Date(string='Published', default=fields.Date.today())
    database = fields.Many2one('rocker.database', string='Datasource',
                               default=lambda self: self.env['rocker.database'].search([('name', '=', 'Odoo')]))
    report_ids = fields.Many2many('rocker.report', 'rocker_report_collection', 'collection_id', 'report_id',
                                  'Collection reports:    (If Excel then check that Sheet names are unique!)',
                                  domain="[('report_type', '=', 'single'),('report_application','=', report_application)]")
    sequence = fields.Integer(string='Sequence', default=10)
    collection_ids = fields.Many2many('rocker.report', 'rocker_report_collection', 'report_id', 'collection_id',
                                      'Report in Collections', domain="[('report_type', '=', 'collection'),('report_application','=', report_application)]")
    column_headings = fields.Char('Column headings', default='Stage; Count', help="Column headings separated with ;")
    select_clause = fields.Text('Select', default=
    """select ns.name, count(*) 
    from note_note nn
    join note_stage_rel nsr on nsr.note_id = nn.id
    join note_stage ns on ns.id = nsr.stage_id 
    group by ns.name, ns.sequence
    order by ns.sequence""")
    sheet_name = fields.Char('Excel Sheet Name', default='Data')
    report_template = fields.Binary('Report template', help="")
    report = fields.Binary('Lastest Report')
    author_id = fields.Many2one('res.users', string='Author', default=lambda self: self.env.user)
    date_executed = fields.Datetime()
    file_name = fields.Char('Report Filename', size=64, default='report.xlsx')
    template_name = fields.Char('Template Filename', size=64, help="")
    perma_link = fields.Char('Permanent link to latest')
    execute_link = fields.Char('Execute & download link')
    interval_number = fields.Integer("Interval", default=1)
    execute_at = fields.Float("Execute at (timezone=UTC)")
    firstcall = fields.Date('First Execution')
    nextcall = fields.Datetime('Next Execution')
    interval_type = fields.Selection(
        [('min', 'min'), ('hour', 'hour'), ('day', 'day'), ('month', 'month')], 'Execute report every', default='day')
    company_id = fields.Many2one('res.company', string='User belonging this company hierarchy can view report')
    _sqldriver = False
    #
    report_application = fields.Selection([('excel', 'Excel Report'), ('powerpoint', 'PowerPoint Report')], 'Report App', default='excel', required=True)
    element = fields.Selection([('table', 'Table'), ('chart', 'Chart')], 'Element type', default='table', required=False)
    legend = fields.Selection([('none', 'None'), ('bottom', 'Bottom'), ('right', 'Right')], 'Legend', default='bottom', required=False)
    elements_per_slide = fields.Selection([('1', '1'), ('2', '2'), ('4', '4'), ('6', '6')], 'Elements per Slide', default='1', required=False)
    slide_title = fields.Char('Slide Title')
    # state = fields.Char('State', store=False, default = 'Testing')
    show_values = fields.Selection([('none', 'None'), ('inside', 'Inside'), ('outside', 'Outside'), ('pros_inside', '% Inside'), ('pros_outside', '% Outside')],'Show values', default='none')
    #python-pptx supports adding charts and modifying existing ones. Most chart types other than 3D types are supported.
    chart_type = fields.Selection([
        #('-4098','3D Area'),
        #('78','3D Stacked Area'),
        #('79','100% Stacked Area'),
        #('60','3D Clustered Bar'),
        #('61','3D Stacked Bar'),
        #('62','3D 100% Stacked Bar'),
        #('-4100','3D Column'),
        #('54','3D Clustered Column'),
        #('55','3D Stacked Column'),
        #('56','3D 100% Stacked Column'),
        #('-4101','3D Line'),
        #('-4102','3D Pie'),
        #('70','Exploded 3D Pie'),
        ('1','Area'),
        ('76','Stacked Area'),
        ('77','100% Stacked Area'),
        ('57','Clustered Bar'),
        #('71','Bar of Pie'),
        ('58','Stacked Bar'),
        ('59','100% Stacked Bar'),
        ('15','Bubble'),
        #('87','Bubble with 3D effects'),
        ('51','Clustered Column'),
        ('52','Stacked Column'),
        ('53','100% Stacked Column'),
        #('102','Clustered Cone Bar'),
        #('103','Stacked Cone Bar'),
        #('104','100% Stacked Cone Bar'),
        #('105','3D Cone Column'),
        #('99','Clustered Cone Column'),
        #('100','Stacked Cone Column'),
        #('101','100% Stacked Cone Column'),
        #('95','Clustered Cylinder Bar'),
        #('96','Stacked Cylinder Bar'),
        #('97','100% Stacked Cylinder Bar'),
        #('98','3D Cylinder Column'),
        #('92','Clustered Cone Column'),
        #('93','Stacked Cone Column'),
        #('94','100% Stacked Cylinder Column'),
        ('-4120','Doughnut'),
        ('80','Exploded Doughnut'),
        ('4','Line'),
        ('65','Line with Markers'),
        ('66','Stacked Line with Markers'),
        ('67','100% Stacked Line with Markers'),
        ('63','Stacked Line'),
        ('64','100% Stacked Line'),
        ('5','Pie'),
        ('99995','Pie  (category from rows)'),
        ('69','Exploded Pie'),
        ('999969','Exploded Pie (category from rows)'),
        #('68','Pie of Pie'),
        #('109','Clustered Pyramid Bar'),
        #('110','Stacked Pyramid Bar'),
        #('111','100% Stacked Pyramid Bar'),
        #('112','3D Pyramid Column'),
        #('106','Clustered Pyramid Column'),
        #('107','Stacked Pyramid Column'),
        #('108','100% Stacked Pyramid Column'),
        ('-4151','Radar'),
        ('82','Filled Radar'),
        ('81','Radar with Data Markers'),
        #('140','Map chart'),
        #('88','High-Low-Close'),
        #('89','Open-High-Low-Close'),
        #('90','Volume-High-Low-Close'),
        #('91','Volume-Open-High-Low-Close'),
        #('83','3D Surface'),
        #('85','Surface (Top View)'),
        #('86','Surface (Top View wireframe)'),
        #('84','3D Surface (wireframe)'),
        ('-4169','Scatter'),
        ('74','Scatter with Lines'),
        ('75','Scatter with Lines and No Data Markers'),
        ('72','Scatter with Smoothed Lines'),
        ('73','Scatter with Smoothed Lines and No Data Markers'),
        ], 'Chart type', required=False)
        # email
    send_by_email = fields.Boolean('Send by Email', default=False)
    email_subject = fields.Char('Subject', default='Rocker Report Notification: [NAME], [DATE]')
    email_to = fields.Char('Email To')
    email_body = fields.Text('Message Body', default='[FILENAME] has been executed at [DATETIME]<p>Cheers<br/>Rocker')
     
    # def write(self, vals):
    #     _logger.debug('Rocker report write')
    #     return super(Report, self).write(vals)

    def export_report(self, context=None):
        if self.active != True:
            raise exceptions.ValidationError('Report is not active or you are not allowed to view it!')
        _logger.info('Rocker reporting / Executing report: ' + self.name)
        if self.report_application == 'excel':
            self.export_xls(self)
        elif self.report_application == 'powerpoint':
            self.export_ppt(self)
        else:
            raise exceptions.ValidationError('Report type?')

        # download
        return {
                'type': 'ir.actions.act_url',
                'name': 'report',
                'url': '/web/content/rocker.report/%s/report/%s?download=true' % (self.id, self.file_name),
            }
 
    def export_ppt(self, context=None):
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        filename = ''
        pp_title = ''
        if not self.file_name:
            self.file_name = 'report.pptx'
        filename = self.file_name.strip()
        if not '.ppt' in filename:
            _logger.debug('Changing name ')
            self.file_name = filename.replace('.xlsx','')
            self.file_name = self.file_name + '.pptx'
            filename = self.file_name.strip()
        odoo_filename = self.file_name.strip()
        # temp_filename = ''
        if not odoo_filename:
            odoo_filename = 'report.pptx'
        if self.slide_title:
            pp_title = self.slide_title.strip()
        else:
            pp_title = '' # case powerpoint slide without title

        output = io.BytesIO()

        #        # take template if exists
        if self.report_template:
            file = base64.b64decode(self.report_template)
            fileobj = tempfile.TemporaryFile('wb+')
            fileobj.write(file)
            fileobj.seek(0)
            prs = Presentation(fileobj)
            _logger.debug('Using template')
        else:
            # generate empty excel
            filename = self.file_name
            _logger.debug('Using empty Powerpoint: ' + filename)
            prs = Presentation()

        if self.report_type == 'single':
            _logger.debug('Single PowerPoint report')
            slide = None
            pp_title = ''
            pp_title = self.slide_title
            collection = False
            cnt_report = 1
            element_written = 1
            page_elements = '1'
            slides_created = 0

            con = self._create_connection(self.database)

            if not pp_title:
                slide = prs.slides.add_slide(prs.slide_layouts[6])   # blanc de blanc
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[5])    # title only
                title = slide.shapes.title
                title.text = pp_title # collection report title set
            slides_created += 1

            # let's fill with data
            if con:
                _logger.debug('Populate single Powerpoint report')
                # one per page
                slide, element_written = self._populate_pp_sql(self, con, prs, slide, page_elements, pp_title, self.select_clause, self.column_headings, cnt_report, element_written, collection)
            else:
                raise exceptions.ValidationError('No DB connection')
        elif self.report_type == 'collection':
            _logger.debug('Executing reports: ')
            _logger.debug(self.report_ids)
            # do loop here
            if not self.elements_per_slide:
                self.elements_per_slide = '1'
            page_elements = self.elements_per_slide
            pp_title = ''
            pp_title = self.slide_title
            slide = None
            cnt_report = 1
            element_written = 1
            slides_created = 0
            collection = True
            for report in self.report_ids:
                con = self._create_connection(report.database)
                if (cnt_report == 1):
                    if not pp_title:
                         slide = prs.slides.add_slide(prs.slide_layouts[6])   # blanc de blanc
                    else:
                        slide = prs.slides.add_slide(prs.slide_layouts[5])    # title only
                        title = slide.shapes.title
                        title.text = pp_title # collection report title set
                    slides_created += 1
                    _logger.debug('Report number ' +  str(cnt_report) + ' on Slide: ' + str(slides_created))
                if con:
                    _logger.debug('PowerPoint Collection report populate')
                    slide, element_written = self._populate_pp_sql(report, con, prs, slide, page_elements, pp_title, report.select_clause, report.column_headings, cnt_report, element_written, collection)
                else:
                    raise exceptions.ValidationError('No DB connection')
                cnt_report += 1
                _logger.debug('Count report sofar: ' + str(cnt_report))
                _logger.debug('Element written sofar: ' + str(element_written))
                if cnt_report > int(page_elements):
                    cnt_report = 1
                if element_written == 1: # we need new slide, next one coming to place 1
                    cnt_report = 1


            # we use temp for saving excel

        prs.save(output)

        # slides ready let's save to database
        datenow = fields.datetime.now()
        _logger.debug('Open file for storing to Odoo: ' + filename)
        # save to report log
        if self.store_history:
            _export_id = self.sudo().env['rocker.archive'].create({'name': self.name, 'date_executed': datenow,
                                                                   'report_file': base64.b64encode(output.getvalue()),
                                                                   'file_name': odoo_filename}, ).id
        # save generated excel
        base_url = self.env['ir.config_parameter'].sudo().get_param('web.base.url')
        _logger.debug('base_url: ' + base_url)
        permlink = base_url + '/web/login?db=' + self._cr.dbname + '&redirect=' + base_url + '/web/content/rocker.report/%s/report/%s?download=true' % (
            self.id, self.file_name)
        _logger.debug('permlink: ' + permlink)
        execlink = self.env['ir.config_parameter'].sudo().get_param(
            'web.base.url') + '/web/login?db=' + self._cr.dbname + '&redirect=' + self.env[
                       'ir.config_parameter'].sudo().get_param(
            'web.base.url') + '/web?#model=rocker.report&view_type=list&menu_id=' + str(
            self.env.ref('rocker_app_v3.rocker_menu').id) + '&action=' + str(
            self.env.ref('rocker_app_v3.rocker_report_execute_request').id) + '&id=' + str(self.id)
        _logger.debug('exec_link: ' + execlink)
        # file.seek(0)
        self.sudo().write({
            'file_name': odoo_filename,
            'report': base64.b64encode(output.getvalue()),
            'date_executed': datenow,
            'perma_link': permlink,
            'execute_link': execlink,
        })

        return True

    def export_xls(self, context=None):
        worksheet = None
        filename = ''
        if not self.file_name:
            self.file_name = 'report.xlsx'
        if not self.sheet_name:
            self.sheet_name = 'Data'
        odoo_filename = self.file_name.strip()
        filename = self.file_name.strip()
        sheetname = self.sheet_name.strip()
        # temp_filename = ''
        if not odoo_filename:
            odoo_filename = 'report.xlsx'
        if not sheetname:
            sheetname = 'Data'

        output = io.BytesIO()

        # take template if exists
        if self.report_template:
            file = base64.b64decode(self.report_template)
            fileobj = tempfile.TemporaryFile('wb+')
            fileobj.write(file)
            fileobj.seek(0)
            # workbook = op.load_workbook(fileobj, data_only=True)
            try:
                workbook = op.load_workbook(fileobj)
            except Exception as e:
                raise exceptions.ValidationError('Error opening template\n\nError text: \n' + str(e))
            _logger.debug('Using template')
        else:
            # generate empty excel
            filename = self.file_name
            _logger.debug('Using empty Excel: ' + filename)
            workbook = op.Workbook()
            sheet = workbook.create_sheet(title=sheetname)
            self._about(workbook)
            # then activate data sheet
            worksheet_names = workbook.sheetnames
            sheet_index = worksheet_names.index(sheetname)
            workbook.active = sheet_index

        if self.report_type == 'single':
            worksheet = self._worksheet(workbook, self.sheet_name)
            _logger.debug('Single report, using sheet: ' + self.sheet_name)
            con = self._create_connection(self.database)
            # let's fill with data
            if con:
                _logger.debug('Populate single report')
                self._populate_sql(con, workbook, worksheet, self.select_clause, self.column_headings)
            else:
                raise exceptions.ValidationError('No DB connection')
        elif self.report_type == 'collection':
            # do loop here
            for report in self.report_ids:
                sheetname = report.sheet_name.strip()
                if not sheetname:
                    raise exceptions.ValidationError('In Collections sheet names are manadatory')
                worksheet = self._worksheet(workbook, sheetname)
                _logger.debug('Collection report, using sheet: ' + report.sheet_name)
                con = self._create_connection(report.database)
                # let's fill with data
                if con:
                    _logger.debug('Collection report populate')
                    self._populate_sql(con, workbook, worksheet, report.select_clause, report.column_headings)
                else:
                    raise exceptions.ValidationError('No DB connection')

        # Refresh all pivot tables
        try:
            pivot = worksheet._pivots[0] # any will do as they share the same cache
            pivot.cache.refreshOnLoad = True
        except:
            _logger.debug('Pivot table refresh problem')

        # save the  workbook
        workbook.save(output)
        workbook.close()
        # workbook ready let's save to database
        datenow = fields.datetime.now()
        _logger.debug('Open file for storing to Odoo: ' + filename)
        # save to report log
        if self.store_history:
            _export_id = self.sudo().env['rocker.archive'].create({'name': self.name, 'date_executed': datenow,
                                                                 'report_file': base64.b64encode(output.getvalue()),
                                                                 'file_name': odoo_filename}, ).id
        # save generated excel
        base_url = self.env['ir.config_parameter'].sudo().get_param('web.base.url')
        permlink = base_url + '/web/login?db=' + self._cr.dbname + '&redirect=' + base_url + '/web/content/rocker.report/%s/report/%s?download=true' % (
        self.id, self.file_name)
        execlink = self.env['ir.config_parameter'].sudo().get_param(
            'web.base.url') + '/web/login?db=' + self._cr.dbname + '&redirect=' + self.env[
                       'ir.config_parameter'].sudo().get_param(
            'web.base.url') + '/web?#model=rocker.report&view_type=list&menu_id=' + str(
            self.env.ref('rocker_app_v3.rocker_menu').id) + '&action=' + str(
            self.env.ref('rocker_app_v3.rocker_report_execute_request').id) + '&id=' + str(self.id)
        # file.seek(0)
        self.sudo().write({
            'file_name': odoo_filename,
            'sheet_name': sheetname,
            'report': base64.b64encode(output.getvalue()),
            'date_executed': datenow,
            'perma_link': permlink,
            'execute_link': execlink,
        })

        return True

    def _worksheet(self, workbook, sheet, context=None):
        _logger.debug('Trying to find sheet: ' + sheet)
        if not sheet:
            sheet = 'Data'

        # find if sheet name exists in template or in workbook
        try:
            worksheet = workbook[sheet]
            _logger.debug('Found worksheet ' + worksheet.title)
        except:
            _logger.debug('Exception sheet creation')
            worksheet = workbook.create_sheet(title=sheet)
            _logger.debug('Created new sheet: ' + worksheet.title)
        return worksheet

    def _about(self, workbook, context=None):
        _logger.debug('Trying to find about sheet: ')
        # we do not create about to template, only to new workbook :-)
        try:
            aboutsheet = workbook['About']
            _logger.debug('Found existing about sheet ')
            aboutsheet.Columns.ClearContents()
        except:
            _logger.debug('About sheet not found, adding new')
            aboutsheet = workbook.create_sheet(title='About')

        aboutsheet.cell(2,2).value = 'Please donate!'
        aboutsheet['C2'].hyperlink = "https://www.paypal.com/cgi-bin/webscr?cmd=_donations&business=DGK3E2CC42EJ4&item_name=for+Rocker+Reporting+application+development&currency_code=EUR&source=url"
        aboutsheet['C2'].value = "Donate with PayPal"
        aboutsheet['C2'].style = "Hyperlink"
        aboutsheet['C2'].hyperlink.tooltip = "Click to Donate"
        aboutsheet.column_dimensions["A"].width = 5
        aboutsheet.column_dimensions["B"].width = 15
        aboutsheet.column_dimensions["C"].width = 100
        for cells in aboutsheet['A2:C2']:
            for cell in cells:
                cell.font = Font(bold=True, size=12)
            break
        aboutsheet.cell(4,2).value = 'Date executed:'
        aboutsheet.cell(4,3).value = fields.datetime.now()
        aboutsheet.cell(4,3).alignment = Alignment(horizontal='left')
        aboutsheet.cell(5,2).value = 'Datasource:'
        aboutsheet.cell(5,3).value = self.name
        if self.report_type == 'collection':
            aboutsheet.cell(6,2).value = 'Type:'
            aboutsheet.cell(6,3).value = 'Collection'
        else:
            aboutsheet.cell(6,2).value = 'SQL:'
            aboutsheet.cell(6,2).alignment = Alignment(wrap_text=True,vertical='top')
            aboutsheet.cell(6,3).value = self.select_clause
            aboutsheet.cell(6,3).alignment = Alignment(wrap_text=True,vertical='top')
        # aboutsheet.Columns.AutoFit()
        # aboutsheet.Rows.AutoFit()

        return True

    def _create_connection(self, database):
        _database_record = self.env['rocker.database'].browse(database.mapped('id'))
        _datasource = _database_record.name
        _driver = _database_record.driver
        _odbcdriver = _database_record.odbcdriver
        _sid = _database_record.database
        _database = _database_record.database
        _host = _database_record.host
        _port = _database_record.port
        _user = _database_record.user
        _password = _database_record.password
        _logger.info('Connecting to ' + _database_record.name)
        con = None
        con = rocker_connection.rocker_connection.create_connection(_database_record)
        if con is not None:
            _logger.info('Database Connect OK')
            return con
        else:
            raise exceptions.ValidationError('Exception, No Database connection')

    def _populate_sql(self, con, workbook, worksheet, sql, headings, context=None):

        for tbl in worksheet.tables:
            _logger.debug('Warning: Sheet ' + worksheet.title + ' has a table ' + tbl + '. If it is overlapping with data area, Excel might be corrupted')

        _headerlist = headings.split(';')
        header = [head.strip() for head in _headerlist]
        cols = len(header)

        first_row= 1
        col = 1
        max_row = 0
        max_in_range = 0
        _logger.debug('Headers')
        for head in header:
            _logger.debug('Col letter:' + get_column_letter(worksheet.cell(row=1,column=col).column))
            try:
                max_row_for = max((c.row for c in worksheet[get_column_letter(worksheet.cell(row=1,column=col).column)] if c.value is not None))
                _logger.debug('max row in col ' + get_column_letter(worksheet.cell(row=1,column=col).column) + ": " + str(max_row_for))
            except:
                max_row_for = 0
            if max_row_for >= max_in_range:
                max_in_range = max_row_for
            worksheet.cell(first_row, col).value = head
            col += 1
        _logger.debug ("Clearing " + str(max_in_range) + " old rows")
        for i in range(2, max_in_range+1):
            for j in range (1,len(header)+1):
                worksheet.cell(row=i,column=j).value = None
        _logger.debug ("Adding data to " + worksheet.title)
        try:
            values = pd.read_sql_query(sql, con)
        except Exception as e:
            raise exceptions.ValidationError('SQL Error!\n\n' + str(e))

        i = 0
        for row in values.values:
            i += 1
            j = 0
            for col in row:
                worksheet.cell(row=i+1,column=j+1).value = col
                j += 1
        # find new max row in data range
        col = 1
        max_in_range = 0
        for head in header:
            max_row_for = max((c.row for c in worksheet[get_column_letter(worksheet.cell(row=1,column=col).column)] if c.value is not None))
            _logger.debug('Found new max row in col ' +str(get_column_letter(worksheet.cell(row=1,column=col).column)) + ": " + str(max_row_for))
            if max_row_for >= max_in_range:
                max_in_range = max_row_for
            col += 1

        # define named range for every col
        col = 1
        for head in header:
            atext = 'OFFSET(' + worksheet.title + '!$' + get_column_letter(worksheet.cell(row=1,column=col).column) + \
                    '$2,,,' + str(max_in_range - 1) + ')'
            _logger.debug("Named range reference: " + atext)
            try:
                my_range = workbook.defined_names[head]
                my_range.attr_text = atext
                _logger.debug('Modified Range: ' + head + " " + my_range.attr_text)
            except:
                _logger.debug('Range does not exist')
                workbook.create_named_range(head,worksheet, '$A$2')
                my_range = workbook.defined_names[head]
                my_range.attr_text = atext
                _logger.debug('Created New Range: ' + head + " " + my_range.attr_text)

            col += 1
        style = TableStyleInfo(name="TableStyleMedium9", showFirstColumn=False,
                               showLastColumn=False, showRowStripes=True, showColumnStripes=False)
        tableref = "A1:{}{}".format(get_column_letter(cols), str(max_in_range))
        tableName = worksheet.title
        if not self.report_template:
            _logger.debug('Create new table')
            _logger.debug(tableref)
            resTable = Table(displayName=tableName, ref=tableref)
            resTable.tableStyleInfo = style
            worksheet.add_table(resTable)
        else:
            _logger.debug('Try to update table(s)')
            for i, tbl in enumerate(worksheet.tables):
                _logger.debug('Sheet ' + worksheet.title + ' has a table ' + tbl )
                if tbl == tableName:
                    _logger.debug('Change ' + tbl )
                    _logger.debug('Original ref: ' + worksheet.tables[tableName].ref)
                    worksheet.tables[tableName].ref = tableref
                    _logger.debug('Changed to: ' + worksheet.tables[tableName].ref)

        # end
        if con is not None:
            con.close()
        return True

    def _populate_pp_sql(self, report, con, prs, slide, page_elements, pp_title, sql, headings, cnt_report, last_element_written, collection, context=None):
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.chart import XL_LEGEND_POSITION
        from pptx.enum.chart import XL_LABEL_POSITION
        from pptx.util import Pt
        from pptx.util import Inches

        self = report
        _headerlist = headings.split(';')
        header = [head.strip() for head in _headerlist]
        cols = len(header)
        slides_created = 1
        element_written = 1
        rows_per_table = 1
        if last_element_written > 1:    # we have created tables and now chart coming
            element_written = last_element_written
        #
        #
        #there are 72 points per inch
        # 5" x 72 points = 360
        # 11.5" x 72 = 828
        # 0.9 x 72 = 65
        # 2 x 72 = 144
        # slide width = 10"
        # slide height 7.5"
        # height - titlespace 5.5
        # element below title y = +2
        font_size = 18

        if self.element == 'chart':
            if self.chart_type == '15':
                from pptx.chart.data import BubbleChartData
                chart_data = BubbleChartData()
            elif self.chart_type in ['-4169','74','75','72','73']:
                from pptx.chart.data import XyChartData
                chart_data = XyChartData()
            else:
                from pptx.chart.data import CategoryChartData
                chart_data = CategoryChartData()
                categories = header   # but first column contains series
                categories.pop(0)
                if self.chart_type not in ['99995', '999969']:   # categories from row data specialities
                    chart_data.categories = categories
                else:
                    categories = []


        # select
        cur = con.cursor()
        try:
            cur.execute(sql)
        except Exception as e:
            raise exceptions.ValidationError('Error in Select clause!\n\n' + str(e))

        records = cur.fetchall()
        i = len(records)

        j = 0
        r = 0
        prev_seriesname = ''
        prev_categoryname = ''
        cnt = 1
        pievalues = []
        # add data rows
        for row in records:

            if self.element == 'table':
                if len(row) != len(header):
                    raise exceptions.ValidationError('Count of headers is not the same as count of data columns\Separate headers with ;\nOr check your SQL')
                if (r >= rows_per_table):
                    # create slide , first slide created already in export_ppt
                    if (element_written == 1):
                        if not pp_title:
                            slide = prs.slides.add_slide(prs.slide_layouts[6])   # blanc de blanc
                        else:
                            slide = prs.slides.add_slide(prs.slide_layouts[5])    # title only
                            title = slide.shapes.title
                            title.text = pp_title #  report title set

                        slides_created = slides_created + 1
                        _logger.debug('New slide again: ' +  str(slides_created))
                    r = 0

                if (r==0):
                    add_title = False
                    if pp_title:
                        add_title = True
                    # find location
                    element_written, rows_per_table, font_size, x, y, cx, cy = self._find_place(page_elements, element_written, add_title)
                    _logger.debug('Table Element written: ' +  str(element_written))
                    # create table
                    try:
                        shape = slide.shapes.add_table(rows_per_table, cols, x, y, cx, cy)
                    except Exception as e:
                        raise exceptions.ValidationError('Error in Table create!\n\n' + str(e))

                    element_written += 1
                    cnt += 1
                    table = shape.table
                    # format table, set font
                    for i in range(rows_per_table):
                        for j in range(cols):
                            cell = table.cell(i, j)
                            cell.text = ' '
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(font_size)

                # set headers
                    c = 0
                    for col in header:
                        cell = table.cell(0, c)
                        cell.text = header[c]
                        run = table.cell(0, c).text_frame.paragraphs[0].runs[0]
                        run.font.size = Pt(font_size)
                        c = c + 1
                    _logger.debug('Table created & formatted')
                    r = 1
                # set data rows
                _logger.debug('Filling Table row: ' +  str(r))
                c = 0
                for col in row:
                    cell = table.cell(r, c)
                    cell.text = str(row[c])
                    run = table.cell(r, c).text_frame.paragraphs[0].runs[0]
                    run.font.size = Pt(font_size)
                    c = c + 1
                r = r + 1
            elif self.element == 'chart':

                if self.chart_type == '15':
                    _logger.debug('Bubble chart 15 ')
                    #                #series_1.add_data_point(0.7, 2.7, 10)
                    seriesname = str(row[0])
                    _logger.debug(seriesname)
                    valuelist = list(row)
                    valuelist.pop(0) # series name taken away
                    _logger.debug('Valuelist: ')
                    _logger.debug(valuelist)
                    #check if we have x & y values
                    if len(valuelist) != 3:
                        raise exceptions.ValidationError('For Bubble chart we need: Series, ValueX, ValueY, ValueSize \nWe must have 3 values (x & y)\nCheck your SQL')
                    # check if rest are numbers only
                    for item in valuelist:
                        if item is not None:
                            if not (isinstance(item, (int, float, complex)) and not isinstance(item, bool)):
                            #except ValueError:
                                raise exceptions.ValidationError('Series, Value, Value, Value \nNow we have string as value\nCheck your SQL')
                    #tuplelist = tuple(valuelist)
                    if not seriesname == prev_seriesname:
                        _logger.debug('Scatter adding series: ' + seriesname)
                        series = chart_data.add_series(seriesname)
                        prev_seriesname = seriesname
                    try:
                        series.add_data_point(valuelist[0],valuelist[1],valuelist[2])
                    except Exception as e:
                        raise exceptions.ValidationError('Error in Bubble!\n\n' + str(e))
                elif self.chart_type in ['-4169','74','75','72','73']:
                    _logger.debug('Chart XY')
                    seriesname = str(row[0])
                    _logger.debug(seriesname)
                    valuelist = list(row)
                    valuelist.pop(0) # series name taken away
                    _logger.debug('Valuelist: ')
                    _logger.debug(valuelist)
                    #check if we have x & y values
                    if len(valuelist) != 2:
                        raise exceptions.ValidationError('For XY diagram (Scatter) we need: Series, ValueX, ValueY \nWe must have 2 values (x & y)\nCheck your SQL')
                   # check if rest are numbers only
                    for item in valuelist:
                        if item is not None:
                            if not (isinstance(item, (int, float, complex)) and not isinstance(item, bool)):
                            #except ValueError:
                                raise exceptions.ValidationError('Series, Value, Value \nNow we have string as value\nCheck your SQL')
                    #tuplelist = tuple(valuelist)
                    if not seriesname == prev_seriesname:
                        _logger.debug('Scatter adding series: ' + seriesname)
                        series = chart_data.add_series(seriesname)
                        prev_seriesname = seriesname
                    try:
                        series.add_data_point(valuelist[0],valuelist[1])
                    except Exception as e:
                        raise exceptions.ValidationError('Error in XY Chart!\n\n' + str(e))
                else:
                    series = str(row[0])
                    _logger.debug(series)
                    valuelist = list(row)
                    valuelist.pop(0) # series name taken away
                    # check if rest are numbers only
                    _logger.debug('Valuelist: ')
                    _logger.debug(valuelist)
                    for item in valuelist:
                        if item is not None:
                            if not (isinstance(item, (int, float, complex)) and not isinstance(item, bool)):
                            #except ValueError:
                                raise exceptions.ValidationError('Series/Category, Value, Value, Value ... \nNow we have string as value\nCheck your SQL')
                    tuplelist = tuple(valuelist)
                    _logger.debug(tuplelist)
                    if self.chart_type in ['99995', '999969']:
                        if len(valuelist) != 1:
                            raise exceptions.ValidationError('For Pie (Category from row data) chart we need: Category, ValueX\nWe must have 2 values (CategoryName & Value)\nCheck your SQL')
                        if not row[0] == prev_categoryname:
                            _logger.debug('Pie category from row ' + str(row[0]))
                            prev_categoryname = row[0]
                            categories.append(row[0])
                            series = 'series1'
                            pievalues.append(row[1])
                        else:
                            raise exceptions.ValidationError('For Pie (Category from row data) every data row must contain unique Category (column 1 = Category Name)')
                    else:
                        try:
                            chart_data.add_series(series,tuplelist)
                        except Exception as e:
                            raise exceptions.ValidationError('Error in Pie Chart!\n\n' + str(e))

        # chart ready lets add
        if self.element == 'chart':
            add_title = False
            if pp_title:
                add_title = True
            # find location
            element_written, rows_per_table, font_size, x, y, cx, cy = self._find_place(page_elements, element_written, add_title)

            chart = int(self.chart_type)
            if chart == 99995:  # AK speciality Pie
                chart = 5
                chart_data.categories = categories
                chart_data.add_series('series1',pievalues)
            if chart == 999969:  # AK speciality Exploded Pie
                chart = 69
                chart_data.categories = categories
                chart_data.add_series('series1',pievalues)
            try:
                graphic_frame = slide.shapes.add_chart(chart, x, y, cx, cy, chart_data)
            except Exception as e:
                raise exceptions.ValidationError('Error in Graphic frame creation!\n\n' + str(e))

            element_written += 1

            chart = graphic_frame.chart
            chart.font.size = Pt(font_size)
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            #shapes2 = slide.shapes
            shapeX = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)   # add black frame
            fill = shapeX.fill
            fill.solid()
            line = shapeX.line
            line.color.rgb = RGBColor(0, 0, 0)  # add black frame
            #fill.fore_color.rgb = RGBColor(0x01, 0x23, 0x45)
            fill.background() # no color inside

            # show values
            #
            # value setting not working
            #  AREA
            # LINE
            # RADAR
            # BUBBLE
            # SCATTER

            if self.chart_type not in ['1','76','77','15','-4120','80','4','65','66','67','63','64','-4151','82',
                                       '81','-4169','75','72','73','-4169','74','75','72']:

                # stacked bars value only inside
                if self.chart_type in ['58','59','52','53',] and self.show_values in ['outside','pros_inside','pros_outside']:
                    self.show_values = 'inside'
                #
                # percentages only for pie charts
                if self.chart_type not in ['5','99995','69','999969',] and self.show_values in ['pros_inside','pros_outside']:
                    self.show_values = 'none'
                #

                if self.show_values == 'inside':
                    chart.plots[0].has_data_labels = True
                    data_labels = chart.plots[0].data_labels
                    data_labels.show_value = True
                    #data_labels.number_format = '0%'
                    data_labels.position = XL_LABEL_POSITION.INSIDE_END
                elif self.show_values == 'outside':
                    chart.plots[0].has_data_labels = True
                    data_labels = chart.plots[0].data_labels
                    data_labels.show_value = True
                    #data_labels.number_format = '0%'
                    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                elif self.show_values == 'pros_inside':
                    chart.plots[0].has_data_labels = True
                    data_labels = chart.plots[0].data_labels
                    data_labels.show_percentage = True
                    data_labels.show_value = False
                    data_labels.number_format = '0%'
                    data_labels.position = XL_LABEL_POSITION.INSIDE_END
                elif self.show_values == 'pros_outside':
                    chart.plots[0].has_data_labels = True
                    data_labels = chart.plots[0].data_labels
                    data_labels.show_value = False
                    data_labels.show_percentage = True
                    data_labels.number_format = '0%'
                    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            else:
                self.show_values = ''
            #legend
            if self.legend == 'bottom':
                chart.has_legend = True
                # XL_LEGEND_POSITION.BOTTOM
                # XL_LEGEND_POSITION.RIGHT
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            elif self.legend == 'right':
                chart.has_legend = True
                # XL_LEGEND_POSITION.BOTTOM
                # XL_LEGEND_POSITION.RIGHT
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.legend.include_in_layout = False
            else:
                chart.has_legend = False
                # XL_LEGEND_POSITION.BOTTOM
                # XL_LEGEND_POSITION.RIGHT
                #chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                #chart.legend.include_in_layout = False

            # chart title
            chart_title = ''
            if self.slide_title:
                chart_title = self.slide_title
            chart.chart_title.text_frame.text = chart_title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(font_size)

        cur.close()

         # end
        if con is not None:
            con.close()
        #return True
        return slide, element_written

    def _find_place(self, page_elements, element_written, pp_title):
        from pptx.util import Pt
        from pptx.util import Inches

        x1, x2, x3, x4, x5, x6 = [0,0,0,0,0,0]
        y1, y2, y3, y4, y5, y6 = [0,0,0,0,0,0]
        cx, cy = [0,0]
        rows_per_table = 0
        font_size = 12
        # pp_title is True or False
        _logger.debug('Page_elements: ' + str(page_elements))
        _logger.debug('Elements_written: ' + str(element_written))

        if page_elements == '1':
            x1, x2, x3, x4, x5, x6 = [Inches(0.5), Inches(0), Inches(0),0,0,0]
            if pp_title:   # slide title, move element down 1 inches, height -1 Inc
                y1, y2, y3, y4, y5, y6 = [Inches(1.5), Inches(0), Inches(0),0,0,0]
                cx, cy = [Inches(9), Inches(5.7)]
            else:
                y1, y2, y3, y4, y5, y6 = [Inches(0.5), Inches(0), Inches(0),0,0,0]
                cx, cy = [Inches(9), Inches(6.7)]
                _logger.debug('No title coordinates set')
            rows_per_table = 13  # header plus 12 data rows
            font_size = 18
        elif page_elements == '2':
            x1, x2, x3, x4, x5, x6 = [Inches(0.5), Inches(5.0), Inches(0),0,0,0]
            if pp_title:   # slide title, move element down 1 inches, height -1 Inc
                y1, y2, y3, y4, y5, y6 = [Inches(1.5), Inches(1.5), Inches(0),0,0,0]
                cx, cy = [Inches(4.4), Inches(5.7)]
            else:
                y1, y2, y3, y4, y5, y6 = [Inches(0.5), Inches(0.5), Inches(0),0,0,0]
                cx, cy = [Inches(4.4), Inches(6.7)]
            rows_per_table = 13
            font_size = 14
        elif page_elements == '4':
            x1, x2, x3, x4, x5, x6 = [Inches(0.5), Inches(5.0), Inches(0.5), Inches(5.0),0,0]
            if pp_title:   # slide title, move element down 1 inches, height -1 Inc
                y1, y2, y3, y4, y5, y6 = [Inches(1.5), Inches(1.5), Inches(4.5), Inches(4.5),0,0]
                cx, cy = [Inches(4.4), Inches(2.67)]
            else:
                y1, y2, y3, y4, y5, y6 = [Inches(0.5), Inches(0.5), Inches(4.0), Inches(4.0),0,0]
                cx, cy = [Inches(4.4), Inches(3.25)]
            rows_per_table = 8 # use font 14
            font_size = 14
        elif page_elements == '6':
            x1, x2, x3, x4, x5, x6 = [Inches(0.5), Inches(3.6), Inches(6.7), Inches(0.5), Inches(3.6), Inches(6.7)]
            if pp_title:   # slide title, move element down 1 inches, height -1 Inc
                y1, y2, y3, y4, y5, y6 = [Inches(1.5), Inches(1.5), Inches(1.5), Inches(4.4), Inches(4.4), Inches(4.4)]
                cx, cy = [Inches(3.0), Inches(2.8)]
            else:
                y1, y2, y3, y4, y5, y6 = [Inches(0.5), Inches(0.5), Inches(0.5), Inches(4.0), Inches(4.0), Inches(4.0)]
                cx, cy = [Inches(3.0), Inches(2.8)]
            rows_per_table = 8
            font_size = 12
        else:
            _logger.debug('Unknown element count')

        # element coordinates
        if page_elements == '1':
            x = x1
            y = y1
        elif page_elements == '2':
            if element_written == 1:
                x = x1
                y = y1
            elif element_written == 2:
                x = x2
                y = y2
        elif page_elements == '4':
            if element_written == 1:
                x = x1
                y = y1
            elif element_written == 2:
                x = x2
                y = y2
            elif element_written == 3:
                x = x3
                y = y3
            elif element_written == 4:
                x = x4
                y = y4
        elif page_elements == '6':
            if element_written == 1:
                x = x1
                y = y1
            elif element_written == 2:
                x = x2
                y = y2
            elif element_written == 3:
                x = x3
                y = y3
            elif element_written == 4:
                x = x4
                y = y4
            elif element_written == 5:
                x = x5
                y = y5
            elif element_written == 6:
                x = x6
                y = y6
        if element_written == int(page_elements):
            element_written = 0
        return element_written, rows_per_table, font_size, x, y, cx, cy

    @api.model
    def _cron_execute_report(self):
        #                                    """)
        _process_reports = self.env.cr.execute(""" SELECT * FROM rocker_report
                     WHERE active = True
                     AND schedule_onoff = True
                     AND  COALESCE(nextcall, firstcall + to_timestamp(COALESCE(execute_at,0) * 60 * 60)::time )  at time zone 'UTC' <= now()
                     """)
        _records = self.env.cr.fetchall()
        for _report in _records:
            self = self.env['rocker.report'].search([('id', '=', _report[0])])
            _logger.info('Cron execute report: ' + self.name)
            self.export_report()

            # email
            if self.send_by_email == True:
                subject = self.email_subject.strip()
                recipients = self.email_to.strip()
                body = self.email_body.strip()
                subject = subject.replace('[NAME]',self.name.strip())
                subject = subject.replace('[FILENAME]',self.file_name.strip())
                subject = subject.replace('[DATE]',str(date.today()))
                subject = subject.replace('[DATETIME]',datetime.now().strftime("%Y-%m-%d, %H:%M"))
                body = body.replace('[NAME]',self.name.strip())
                body = body.replace('[FILENAME]',self.file_name.strip())
                body = body.replace('[DATE]',str(date.today()))
                body = body.replace('[DATETIME]',datetime.now().strftime("%Y-%m-%d, %H:%M"))

                recipients = ''
                recipients = self.email_to.strip()

                 # now the sender is OdooBot
                template_obj = self.env['mail.mail']
                #template_data = {
                #    'subject': 'Rocker Report Notification: ' + self.name + ' is ready',
                #    'body_html': message_body,
                #    'email_from': sender,
                #    'email_to': ", ".join(recipients),
                #}
                template_data = {
                    'subject': subject,
                    'body_html': body,
                    'email_to': recipients,
                    'auto_delete': True,
                }
                template_id = template_obj.create(template_data)
                attach_obj = self.env['ir.attachment']
                attachment_ids = []
                # filename & datas_fname only for Odoo 12
                #attach_data = {
                #    'name': self.file_name,
                #    'filename': self.file_name,
                #    'datas': self.report,
                #    'datas_fname': self.file_name,
                #    'res_model': 'ir.ui.view',
                #}
                attach_data = {
                    'name': self.file_name,
                    'datas': self.report,
                    'res_model': 'ir.ui.view',
                }
                attach_id = attach_obj.create(attach_data)
                attachment_ids.append(attach_id.id)
                if attachment_ids:
                    template_id.write({'attachment_ids': [(6, 0, attachment_ids)]})

                _logger.debug('Subject ' + subject)
                _logger.debug('Email To ' + recipients)
                _logger.debug('Body ' + body)

                template_id.send()

            # define next run for the report
            if self.interval_type == 'min':
                _logger.debug('Minute Intervall')
                nextd = datetime.now()
                next = datetime.now().strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Executed at: ' + next)
                nextd = datetime.now() + timedelta(minutes=self.interval_number)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Next Execution at: ' + next)
                self.nextcall = nextd
            elif self.interval_type == 'hour':
                _logger.debug('Hour Intervall')
                nextd = datetime.now()
                next = datetime.now().strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Executed at: ' + next)
                nextd = datetime.now() + timedelta(hours=self.interval_number)
                nextd = nextd.replace(minute=0, second=0)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                # take minutes from execute_at field, we don't want to flush exec time
                minutes = self.execute_at * 60
                hours, minutes = divmod(minutes, 60)
                nextd = nextd + timedelta(minutes=minutes)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Next Execution at: ' + next)
                self.nextcall = nextd

            elif self.interval_type == 'day':
                _logger.debug('Define Day Intervall')
                nextd = datetime.now()
                next = datetime.now().strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Executed at: ' + next)
                nextd = datetime.now() + timedelta(days=self.interval_number)
                nextd = nextd.replace(hour=0, minute=0, second=0)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                # take minutes from execute_at field, we don't want to flush exec time
                minutes = self.execute_at * 60
                hours, minutes = divmod(minutes, 60)
                nextd = nextd + timedelta(hours=hours)
                nextd = nextd + timedelta(minutes=minutes)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Next Execution at: ' + next)
                self.nextcall = nextd

            elif self.interval_type == 'month':
                _logger.debug('Define Month Intervall')
                nextd = datetime.now()
                next = datetime.now().strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Executed at: ' + next)
                from dateutil.relativedelta import relativedelta
                nextd = datetime.now() + relativedelta(months=self.interval_number)
                nextd = nextd.replace(hour=0, minute=0, second=0)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                # take minutes from execute_at field, we don't want to flush exec time
                minutes = self.execute_at * 60
                hours, minutes = divmod(minutes, 60)
                nextd = nextd + timedelta(hours=hours)
                nextd = nextd + timedelta(minutes=minutes)
                next = nextd.strftime("%Y-%m-%d, %H:%M")
                _logger.debug('Next Execution at: ' + next)
                self.nextcall = nextd

            else:
                _logger.debug('Unknown intervall')

        _logger.debug('Nothing to do...boooring!')

    @api.model
    def _execute_xls(self, context=None):
        report_id = dict(self._context.get('params', {})).get('id')
        self = self.env['rocker.report'].search([('id', '=', report_id)])
        self.export_report()
        _logger.debug('Base url: ' + self.env['ir.config_parameter'].sudo().get_param('web.base.url'))
        return {
            'type': 'ir.actions.act_url',
            'name': 'report',
            'url': '/web/content/rocker.report/%s/report/%s?download=true' % (self.id, self.file_name)
        }

    @api.model
    def _testexcel(self):
        _logger.debug('Starting test')
        output = io.BytesIO()
        try:
            wb = op.Workbook()
            sheet = wb.create_sheet(title="Data")
            sheet["A1"] = "This is a template!"
            _logger.debug('Save to output ')
            wb.save(output)
            wb.close()
            _logger.debug('Excel quit')
        except Exception as e:
            raise exceptions.ValidationError(
                'Excel test\n\nTried to create files to memory\n\n' + str(e))
        context = {}
        context['message'] = "Excel worksheet creation seems to work!"
        title = 'Success'
        view = self.env.ref('rocker_app_v3.rocker_popup_wizard')
        view_id = False
        return {
            'name': title,
            'type': 'ir.actions.act_window',
            'view_type': 'form',
            'view_mode': 'form',
            'res_model': 'rocker.popup.wizard',
            'views': [(view.id, 'form')],
            'view_id': view.id,
            'target': 'new',
            'context': context,
        }

